"""Extract plain text from every TroySD document into _text/ mirror tree.

Output layout (corpus root from TSD_BOE_ROOT env var, default <repo>/data/tsd-boe-data):
  <root>/_text/<meeting_folder>/<original_filename>.txt

Legacy formats (.doc, .ppt) go through the macOS `textutil` and LibreOffice
`soffice` converters. Images (.png/.jpg) and PDFs with no usable text layer go
through OCR (`ocrmypdf`/`tesseract`).

An extension is a hint, not a fact: files are sniffed by magic bytes when the
suffix has no extractor, which is how a board packet saved as `.tmp` was found.
A PDF that yields text is not necessarily a PDF that yields *usable* text — a
subset font with no ToUnicode CMap decodes to glyph ids like `/0/1/2/i255`, and
a scanned page yields a few bytes of header. Both are re-run through OCR and the
better result wins. See `_needs_ocr()`.

Still skipped and recorded in _text/_skipped.txt: anything with no extractor and
no recognised magic. See extract_legacy.py for a Windows COM-based fallback.
"""
from __future__ import annotations

import os
import re
import shutil
import subprocess
import sys
import tempfile
import time
import zipfile
from pathlib import Path

from pypdf import PdfReader
from pypdf.errors import PdfReadError
import pdfplumber
import docx
from pptx import Presentation
import openpyxl
from striprtf.striprtf import rtf_to_text

ROOT = Path(os.environ.get("TSD_BOE_ROOT") or
        Path(__file__).resolve().parent / "data" / "tsd-boe-data")
TEXT_ROOT = ROOT / "_text"


# A leading digit split off from its own thousands group, e.g. "1 23,879,792".
# The lookbehind keeps genuinely separate columns ("31,855,092 34,278,115") out.
SPLIT_NUM = re.compile(r"(?<![\d,])\d{1,2} \d{2,3},")

# Meeting-folder date, e.g. "2019-06-17_Regular Meeting ...".
MEETING_DATE_RE = re.compile(r"(\d{4}-\d{2}-\d{2})_")
# Full-meeting packet filenames: "061620Mtg", "111114RegMtg", "020717SpMtg".
PACKET_NAME_RE = re.compile(r"^W?\d{6}\s*(reg|sp|org|wksh|wksp|closed)?\s*"
                            r"(mtg|bdmtg|meeting)", re.I)
# Exclusions from the pdfplumber pass — see _worth_reordering().
REORDER_AFTER = os.environ.get("TSD_REORDER_AFTER", "2020-01-01")
MAX_REORDER_BYTES = int(float(os.environ.get("TSD_MAX_REORDER_MB", "15")) * 1e6)
REORDER_PACKETS = os.environ.get("TSD_REORDER_PACKETS", "") not in ("", "0")


def _pypdf_pages(p: Path) -> list[str]:
    try:
        r = PdfReader(str(p), strict=False)
    except (PdfReadError, Exception):
        return []
    out = []
    for page in r.pages:
        try:
            out.append(page.extract_text() or "")
        except Exception:
            out.append("")
    return out


def _plumber_pages(p: Path) -> list[str]:
    try:
        with pdfplumber.open(str(p)) as pdf:
            out = []
            for page in pdf.pages:
                try:
                    out.append(page.extract_text() or "")
                except Exception:
                    out.append("")
            return out
    except Exception:
        return []


def _worth_reordering(p: Path) -> bool:
    """Whether a pdfplumber pass is worth its cost for this file.

    The reorder fixes one artifact: pypdf emitting a heading after the table it
    labels, in documents built from headed tables — financial statements, budget
    books, ACFRs. pdfplumber's cost scales with page count, so it is ruinous on
    the few very long documents that have no such structure. Three exclusions,
    each measured against this corpus:

    * Full-meeting packets, by filename (``061620Mtg.pdf``, ``111114Mtg.pdf``).
      202 PDFs. These are whole meetings bound into one file — hundreds of pages
      of mixed content, no consistent heading/table pairing to repair. Matching
      on name rather than date matters: 2020 is a transition year that still
      contains packets, and one of them held extraction at 100% CPU for four
      minutes on a single file.
    * Anything over MAX_REORDER_BYTES. Catches the image-heavy presentations
      (69 MB, 44 MB, 31 MB) whose slides have no tables worth reordering. Chosen
      so every document that does need it survives: the largest budget book is
      4.8 MB and the largest ACFR 12.8 MB.
    * Meetings before REORDER_AFTER, the era deferred from re-summarization.

    Set TSD_REORDER_AFTER=0000-00-00 and TSD_MAX_REORDER_MB=0 to reorder all.
    Set TSD_REORDER_PACKETS=1 to include full-meeting packets despite the cost;
    measured on the pre-2020 era, the reorder moves 54-73% of a packet's lines
    and fixes real damage (agenda footers emitted ahead of the agenda itself),
    so the exclusion is a speed trade-off rather than a correctness one.
    """
    if PACKET_NAME_RE.match(p.stem) and not REORDER_PACKETS:
        return False
    try:
        if MAX_REORDER_BYTES and p.stat().st_size > MAX_REORDER_BYTES:
            return False
    except OSError:
        pass
    m = MEETING_DATE_RE.match(p.parent.name)
    return not (m and m.group(1) < REORDER_AFTER)


def _key(line: str) -> str:
    """Whitespace-insensitive key for matching a line across the two extractors.

    pdfplumber splits digits off their thousands groups ("1 23,879,792") and
    spaces columns differently from pypdf, but the non-space characters agree.
    """
    return "".join(line.split())


def _reorder(pypdf_txt: str, plumber_txt: str) -> str:
    """Re-emit pypdf's lines in the reading order pdfplumber observed.

    pypdf emits runs in content-stream order; on these documents that scatters
    a page's heading block after its own table, so a sequential reader
    attributes each table to the wrong heading. pdfplumber reads in visual
    order but corrupts figures, so it can supply the order and nothing else.

    Match each pdfplumber line to a pypdf line on non-space characters, emit
    the pypdf originals in pdfplumber's sequence, then append any pypdf lines
    pdfplumber never produced so no text is dropped. Returns pypdf's text
    unchanged when the two sides agree too poorly to trust the alignment.
    """
    ylines = [l for l in pypdf_txt.split("\n")]
    plines = [l for l in plumber_txt.split("\n") if l.strip()]
    if not plines or not ylines:
        return pypdf_txt

    buckets: dict[str, list[int]] = {}
    for i, l in enumerate(ylines):
        if l.strip():
            buckets.setdefault(_key(l), []).append(i)

    order, used = [], set()
    for pl in plines:
        idxs = buckets.get(_key(pl))
        if idxs:
            for i in idxs:
                if i not in used:
                    order.append(i)
                    used.add(i)
                    break

    # Too few matches means the alignment is guesswork — keep pypdf as-is.
    if len(order) < max(3, int(0.6 * len(plines))):
        return pypdf_txt

    order += [i for i in range(len(ylines)) if i not in used and ylines[i].strip()]
    return "\n".join(ylines[i] for i in order)


def text_pdf(p: Path) -> str:
    """Extract a PDF's text with pypdf's characters and pdfplumber's ordering.

    Neither library is correct alone on this corpus. pypdf preserves figures
    exactly but emits page content out of reading order; pdfplumber reads in
    visual order but splits digits off their thousands groups, which silently
    corrupts every dollar amount. Since the corpus exists to be queried for
    dollar amounts, the characters have to come from pypdf.

    So: extract with both, and use pdfplumber only to decide where each page
    starts. pdfplumber is skipped entirely when pypdf's output already looks
    like reading order, and falls back cleanly whenever either side fails.
    """
    ypages = _pypdf_pages(p)
    ytxt = "\n".join(ypages).strip()
    if not ytxt:
        # pypdf found nothing — take pdfplumber's text even with its defects
        return "\n".join(_plumber_pages(p)).strip()

    if not _worth_reordering(p):
        return ytxt

    ppages = _plumber_pages(p)
    if len(ppages) != len(ypages):
        return ytxt                          # can't align pages; trust pypdf

    fixed = [_reorder(y, q) for y, q in zip(ypages, ppages)]
    return "\n".join(fixed).strip()


def text_docx(p: Path) -> str:
    try:
        d = docx.Document(str(p))
    except Exception:
        return ""
    parts = [par.text for par in d.paragraphs]
    for tbl in d.tables:
        for row in tbl.rows:
            for cell in row.cells:
                parts.append(cell.text)
    return "\n".join(parts).strip()


def text_pptx(p: Path) -> str:
    try:
        prs = Presentation(str(p))
    except Exception:
        return ""
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for par in shape.text_frame.paragraphs:
                    parts.append("".join(r.text for r in par.runs))
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
    return "\n".join(parts).strip()


def text_xlsx(p: Path) -> str:
    try:
        wb = openpyxl.load_workbook(str(p), data_only=True, read_only=True)
    except Exception:
        return ""
    parts = []
    for ws in wb.worksheets:
        parts.append(f"=== Sheet: {ws.title} ===")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(cells):
                parts.append("\t".join(cells))
    wb.close()
    return "\n".join(parts).strip()


def text_rtf(p: Path) -> str:
    try:
        return rtf_to_text(p.read_text(encoding="utf-8", errors="replace")).strip()
    except Exception:
        return ""


# --- OCR and legacy-format converters -------------------------------------
# Every binary here ships with the repo's documented toolchain (homebrew
# tesseract/ocrmypdf/libreoffice, plus macOS textutil). A missing one degrades
# to "" and the file lands in _skipped.txt rather than taking the run down.

# Glyph-id runs from a subset font with no ToUnicode CMap: "/0/1/2/i255/5".
CID_RUN = re.compile(r"/\d+/\d+/\d+")
OCR_MIN_CHARS_PER_PAGE = 120   # below this a PDF page is assumed to be an image
OCR_BYTES_PER_CHAR = 3000      # source bytes per extracted char; above this it is pictures


def _have(binary: str) -> bool:
    return shutil.which(binary) is not None


def _is_glyph_garbage(txt: str) -> bool:
    """Text that decoded to font glyph ids rather than characters."""
    return bool(txt) and len(CID_RUN.findall(txt)) * 12 > len(txt) * 0.25


def _needs_ocr(txt: str, pages: int, src_bytes: int = 0) -> bool:
    """True when a PDF 'extracted' but the result is not usable text.

    Three failure modes look like success to every other check in this pipeline:
    a scanned page returns a few bytes of header, a subset font with no
    ToUnicode CMap returns glyph ids that are the right *length* but carry no
    words, and a page of scanned tables returns its caption and nothing else.
    Figure validation cannot see any of them, because none asserts a figure.

    The per-page floor alone is not enough. `MHSAA resolution 23-24.pdf` is
    2.3MB and yielded 302 characters on one page -- comfortably over a 120-char
    floor, and still an image. What gives it away is the ratio: real text costs
    a few hundred source bytes per character, a scan costs thousands.
    """
    if not txt:
        return True
    if _is_glyph_garbage(txt):
        return True
    n = len(txt.strip())
    if n < OCR_MIN_CHARS_PER_PAGE * max(1, pages):
        return True
    return src_bytes > 100_000 and src_bytes / max(1, n) > OCR_BYTES_PER_CHAR


def ocr_pdf(p: Path) -> str:
    if not _have("ocrmypdf"):
        return ""
    with tempfile.TemporaryDirectory() as td:
        side = Path(td) / "out.txt"
        try:
            subprocess.run(["ocrmypdf", "--force-ocr", "--output-type", "none",
                            "--sidecar", str(side), str(p), os.devnull],
                           capture_output=True, timeout=1800, check=False)
            return side.read_text(encoding="utf-8", errors="replace").strip() if side.exists() else ""
        except Exception:
            return ""


def text_image(p: Path) -> str:
    if not _have("tesseract"):
        return ""
    with tempfile.TemporaryDirectory() as td:
        stem = Path(td) / "o"
        try:
            subprocess.run(["tesseract", str(p), str(stem)],
                           capture_output=True, timeout=300, check=False)
            out = stem.with_suffix(".txt")
            return out.read_text(encoding="utf-8", errors="replace").strip() if out.exists() else ""
        except Exception:
            return ""


def text_doc(p: Path) -> str:
    """Legacy Word. macOS textutil first -- it is instant and lossless enough --
    then soffice, which also covers non-Darwin hosts."""
    if _have("textutil"):
        try:
            r = subprocess.run(["textutil", "-convert", "txt", "-stdout", str(p)],
                               capture_output=True, timeout=120)
            got = r.stdout.decode("utf-8", errors="replace").strip()
            if got:
                return got
        except Exception:
            pass
    return _soffice_text(p, "txt")


def text_ppt(p: Path) -> str:
    """Legacy PowerPoint: convert to .pptx and reuse the existing reader."""
    conv = _soffice_convert(p, "pptx")
    return text_pptx(conv) if conv else ""


def _soffice_convert(p: Path, fmt: str):
    if not _have("soffice"):
        return None
    td = tempfile.mkdtemp()
    try:
        subprocess.run(["soffice", "--headless", "--convert-to", fmt,
                        "--outdir", td, str(p)],
                       capture_output=True, timeout=600, check=False)
        hits = list(Path(td).glob(f"*.{fmt}"))
        return hits[0] if hits else None
    except Exception:
        return None


def _soffice_text(p: Path, fmt: str) -> str:
    conv = _soffice_convert(p, fmt)
    if not conv:
        return ""
    try:
        return conv.read_text(encoding="utf-8", errors="replace").strip()
    except Exception:
        return ""


def text_docx_media(p: Path) -> str:
    """A .docx whose body is empty but which wraps one or more images -- OCR them.
    Three board documents are nothing but a screenshot in a Word wrapper."""
    try:
        z = zipfile.ZipFile(p)
    except Exception:
        return ""
    media = [n for n in z.namelist() if n.startswith("word/media/")]
    if not media:
        return ""
    parts = []
    with tempfile.TemporaryDirectory() as td:
        for name in media:
            dst = Path(td) / Path(name).name
            try:
                dst.write_bytes(z.read(name))
            except Exception:
                continue
            got = text_image(dst)
            if got:
                parts.append(got)
    return "\n".join(parts).strip()


def office_ocr(p: Path) -> str:
    """Last resort for an Office file whose slides/pages are just images.

    A 16MB .pptx that yields 967 characters is a deck of screenshots, and
    python-pptx can only read the text boxes that are not there. Rendering the
    whole file to PDF through soffice and OCR-ing that recovers what a reader
    actually sees."""
    conv = _soffice_convert(p, "pdf")
    return ocr_pdf(conv) if conv else ""


# Magic bytes -> the extension we should have been given.
MAGIC = [(b"%PDF", ".pdf"), (b"\x89PNG", ".png"), (b"\xff\xd8\xff", ".jpg"),
         (b"\xd0\xcf\x11\xe0", ".doc")]


def sniff(p: Path):
    try:
        head = p.open("rb").read(8)
    except Exception:
        return None
    for sig, ext in MAGIC:
        if head.startswith(sig):
            return ext
    return None


EXTRACT = {
    ".pdf": text_pdf,
    ".docx": text_docx,
    ".pptx": text_pptx,
    ".xlsx": text_xlsx,
    ".rtf": text_rtf,
    ".doc": text_doc,
    ".ppt": text_ppt,
    ".png": text_image,
    ".jpg": text_image,
    ".jpeg": text_image,
}


def main():
    # --recheck re-opens files that already have output and re-tests it with
    # _needs_ocr(). Without it the existing-output shortcut below means a batch
    # that extracted to glyph ids or to a 12-byte header stays that way forever:
    # the file is present and non-empty, so every later run skips it.
    recheck = "--recheck" in sys.argv
    TEXT_ROOT.mkdir(parents=True, exist_ok=True)
    files = sorted(p for p in ROOT.rglob("*")
                   if p.is_file()
                   and not p.name.startswith("_")
                   and p.suffix.lower() != ".py"
                   and TEXT_ROOT not in p.parents)

    skipped_log = (TEXT_ROOT / "_skipped.txt").open("w", encoding="utf-8")
    summary = {"ok": 0, "empty": 0, "skipped": 0, "error": 0, "ocr": 0, "sniffed": 0, "rechecked": 0}
    t0 = time.time()
    n = len(files)

    for i, p in enumerate(files, 1):
        ext = p.suffix.lower()
        rel = p.relative_to(ROOT)
        out = TEXT_ROOT / rel.with_suffix(rel.suffix + ".txt")
        out.parent.mkdir(parents=True, exist_ok=True)
        if out.exists() and out.stat().st_size > 0:
            if not recheck:
                summary["ok"] += 1
                continue
            prior = out.read_text(encoding="utf-8", errors="replace")
            body = prior.split("=" * 30)[-1] if "=" * 30 in prior else prior
            if ext in (".docx", ".pptx", ".doc", ".ppt", ".rtf"):
                if len(body.strip()) >= 1000 or p.stat().st_size <= 150_000:
                    summary["ok"] += 1
                    continue
                summary["rechecked"] += 1
                # fall through: re-extract, then the office-OCR branch below
            elif ext != ".pdf":
                summary["ok"] += 1
                continue
            try:
                pages = len(PdfReader(str(p)).pages)
            except Exception:
                pages = 1
            if not _needs_ocr(body, pages, p.stat().st_size):
                summary["ok"] += 1
                continue
            summary["rechecked"] += 1
        fn = EXTRACT.get(ext)
        if fn is None:
            # The suffix is a hint, not a fact. A whole board workshop was
            # sitting in the corpus as `101811WkspMtg.tmp` -- a PDF that the
            # extension filter walked straight past.
            real = sniff(p)
            fn = EXTRACT.get(real) if real else None
            if fn is None:
                skipped_log.write(f"{rel}\t(no extractor for {ext})\n")
                summary["skipped"] += 1
                continue
            ext = real
            summary["sniffed"] += 1
        try:
            txt = fn(p)
        except Exception as e:
            skipped_log.write(f"{rel}\tERROR: {e}\n")
            summary["error"] += 1
            continue

        # A PDF that returned *something* has not necessarily returned usable
        # text. Re-run the thin and the glyph-garbled through OCR and keep
        # whichever is longer -- OCR is not always better, and on a genuinely
        # near-empty page it comes back shorter, which is the honest answer.
        if ext == ".pdf":
            try:
                pages = len(PdfReader(str(p)).pages)
            except Exception:
                pages = 1
            if _needs_ocr(txt, pages, p.stat().st_size):
                got = ocr_pdf(p)
                # "Keep whichever is longer" is right for a thin scan, where OCR
                # either finds words or honestly finds none. It is exactly WRONG
                # for glyph garbage: a subset font with no ToUnicode CMap decodes
                # to tens of thousands of characters of "/16/17/18/i255", which
                # beats any real OCR result on length and is worth nothing. When
                # the incumbent text is garbage, any real text wins.
                if got and (_is_glyph_garbage(txt) or len(got) > len(txt or "")):
                    txt = got
                    summary["ocr"] += 1
        elif ext in (".docx", ".pptx", ".doc", ".ppt", ".rtf"):
            # Same trap, different container: the text layer is thin because the
            # content is pictures. Try embedded media first (cheap), then render
            # the whole file and OCR it. Keep whichever is longer.
            if not txt and ext == ".docx":
                got = text_docx_media(p)
                if got:
                    txt = got
                    summary["ocr"] += 1
            if len((txt or "").strip()) < 1000 and p.stat().st_size > 150_000:
                got = office_ocr(p)
                if len(got) > len(txt or ""):
                    txt = got
                    summary["ocr"] += 1

        if not txt:
            skipped_log.write(f"{rel}\t(empty extraction)\n")
            summary["empty"] += 1
            continue
        out.write_text(txt, encoding="utf-8")
        summary["ok"] += 1
        if i % 50 == 0 or i == n:
            elapsed = time.time() - t0
            print(f"[{i:>4}/{n}] {elapsed:6.1f}s  {summary}", flush=True)

    skipped_log.close()
    print(f"DONE {summary}")


if __name__ == "__main__":
    sys.exit(main() or 0)
