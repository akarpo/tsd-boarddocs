"""Estimate token count for the downloaded TroySD documents.

For each supported file type, extract plain text and count tokens with
tiktoken cl100k_base (GPT-4 tokenizer) — a reasonable proxy for Claude's
tokenizer (Claude tokens for English text are typically within ~5-10%).

Output: per-type and grand totals to stdout, plus _tokens_per_file.csv
under the corpus root (TSD_BOE_ROOT env var, default ~/tsd-boe-data).
"""
from __future__ import annotations

import csv
import io
import os
import sys
import time
from pathlib import Path

import tiktoken
from pypdf import PdfReader
from pypdf.errors import PdfReadError
import docx
from pptx import Presentation
import openpyxl
from striprtf.striprtf import rtf_to_text

ROOT = Path(os.environ.get("TSD_BOE_ROOT") or Path.home() / "tsd-boe-data")
ENC = tiktoken.get_encoding("cl100k_base")


def text_pdf(p: Path) -> str:
    parts = []
    try:
        r = PdfReader(str(p), strict=False)
        for page in r.pages:
            try:
                parts.append(page.extract_text() or "")
            except Exception:
                pass
    except (PdfReadError, Exception):
        pass
    return "\n".join(parts)


def text_docx(p: Path) -> str:
    try:
        d = docx.Document(str(p))
    except Exception:
        return ""
    parts = [par.text for par in d.paragraphs]
    for tbl in d.tables:
        for row in tbl.rows:
            for cell in row.cells:
                parts.append(cell.text)
    return "\n".join(parts)


def text_pptx(p: Path) -> str:
    try:
        prs = Presentation(str(p))
    except Exception:
        return ""
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for par in shape.text_frame.paragraphs:
                    parts.append("".join(r.text for r in par.runs))
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
    return "\n".join(parts)


def text_xlsx(p: Path) -> str:
    try:
        wb = openpyxl.load_workbook(str(p), data_only=True, read_only=True)
    except Exception:
        return ""
    parts = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            for v in row:
                if v is not None:
                    parts.append(str(v))
    wb.close()
    return "\n".join(parts)


def text_rtf(p: Path) -> str:
    try:
        return rtf_to_text(p.read_text(encoding="utf-8", errors="replace"))
    except Exception:
        return ""


EXTRACT = {
    ".pdf": text_pdf,
    ".docx": text_docx,
    ".pptx": text_pptx,
    ".xlsx": text_xlsx,
    ".rtf": text_rtf,
}


def main():
    files = [p for p in ROOT.rglob("*") if p.is_file()
             and not p.name.startswith("_")
             and p.suffix.lower() != ".py"]
    by_ext = {}
    csv_path = ROOT / "_tokens_per_file.csv"
    csvf = csv_path.open("w", encoding="utf-8", newline="")
    w = csv.writer(csvf)
    w.writerow(["path", "ext", "size_bytes", "chars", "tokens", "extracted"])

    t0 = time.time()
    n = len(files)
    for i, p in enumerate(files, 1):
        ext = p.suffix.lower()
        size = p.stat().st_size
        fn = EXTRACT.get(ext)
        if fn is None:
            # legacy/unsupported — record metadata only
            w.writerow([str(p), ext, size, 0, 0, 0])
            d = by_ext.setdefault(ext, {"files": 0, "bytes": 0, "chars": 0,
                                        "tokens": 0, "extracted": 0})
            d["files"] += 1
            d["bytes"] += size
            continue
        try:
            txt = fn(p)
        except Exception:
            txt = ""
        chars = len(txt)
        tokens = len(ENC.encode(txt, disallowed_special=())) if txt else 0
        d = by_ext.setdefault(ext, {"files": 0, "bytes": 0, "chars": 0,
                                    "tokens": 0, "extracted": 0})
        d["files"] += 1
        d["bytes"] += size
        d["chars"] += chars
        d["tokens"] += tokens
        d["extracted"] += 1 if chars else 0
        w.writerow([str(p), ext, size, chars, tokens, 1 if chars else 0])
        if i % 50 == 0 or i == n:
            elapsed = time.time() - t0
            print(f"  [{i}/{n}] elapsed={elapsed:6.1f}s  last={p.name[:60]}",
                  flush=True)
    csvf.close()

    # Estimate legacy .doc / .ppt by applying their modern-cousin token-per-byte rate
    docx_d = by_ext.get(".docx", {"tokens": 0, "bytes": 1})
    pptx_d = by_ext.get(".pptx", {"tokens": 0, "bytes": 1})
    docx_rate = docx_d["tokens"] / max(docx_d["bytes"], 1)
    pptx_rate = pptx_d["tokens"] / max(pptx_d["bytes"], 1)
    legacy_estimate = 0
    for ext, rate in ((".doc", docx_rate), (".ppt", pptx_rate)):
        d = by_ext.get(ext)
        if d:
            est = int(d["bytes"] * rate)
            d["tokens_estimate"] = est
            legacy_estimate += est

    print()
    print(f"{'ext':<8}{'files':>8}{'extracted':>11}{'bytes':>15}"
          f"{'chars':>15}{'tokens':>15}{'tok/file':>10}")
    print("-" * 82)
    grand_tokens = 0
    for ext in sorted(by_ext):
        d = by_ext[ext]
        tok = d["tokens"] if d["tokens"] else d.get("tokens_estimate", 0)
        avg = tok // d["files"] if d["files"] else 0
        marker = "*" if d.get("tokens_estimate") else " "
        print(f"{ext:<8}{d['files']:>8}{d['extracted']:>11}"
              f"{d['bytes']:>15,}{d['chars']:>15,}{tok:>14,}{marker}{avg:>9,}")
        grand_tokens += tok
    print("-" * 82)
    print(f"GRAND TOTAL tokens: {grand_tokens:,}")
    print(f"  (* = estimated from sibling-format token-per-byte rate)")


if __name__ == "__main__":
    main()
